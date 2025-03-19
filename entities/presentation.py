from datetime import datetime
from typing import TypedDict
from pptx import Presentation as PPTX
from pptx.util import Cm
from utils import create_supplier_ref_list, text_frame_paragraph

from entities.entities import Contact, Representative, Client


class Presentation:
    prs = PPTX("./plantillas/cotizacion.pptm")
    title_slide_layout = prs.slide_layouts[6]

    def __init__(self, references_quantity: int) -> None:
        self.ref_q = references_quantity

    def set_slides(self, contact: Contact):
        for i in range(0, self.ref_q + 1):
            self.prs.slides.add_slide(self.title_slide_layout)
            self.prs.slides[i].shapes.add_picture(
                "./images/logo.jpeg",
                left=Cm(1),
                top=Cm(0.5),
                width=Cm(8.9),
                height=Cm(1.7),
            )
            footer_textbox = self.prs.slides[i].shapes.add_textbox(
                left=Cm(0.5), top=Cm(22.8), width=Cm(18), height=Cm(1)
            )
            tf_footer = footer_textbox.text_frame
            text_frame_paragraph(
                tf_footer,
                f"{contact['address']} {contact['phone']} {contact['email']} {contact['web']}",
                7,
                False,
                True,
            )

    def add_header(self, representative: Representative, consecutive: int):
        today = datetime.now()
        txBox = self.prs.slides[0].shapes.add_textbox(
            left=Cm(12), top=Cm(-0.5), width=Cm(6.6), height=Cm(6)
        )
        tf = txBox.text_frame
        text_frame_paragraph(
            tf, f"{today.day} {today.strftime('%B')} de {today.year}", 14
        )
        text_frame_paragraph(tf, f"Cot N°{consecutive}", 14)
        text_frame_paragraph(tf, "", 11)
        text_frame_paragraph(tf, "Asesor Comercial", 11)
        text_frame_paragraph(
            tf,
            f"{representative['name']} {representative['phone']} {representative['email']}",
            11,
        )

    def add_client_name(self, client: Client):
        header = self.prs.slides[0].shapes.add_textbox(
            left=Cm(1), top=Cm(1.8), width=Cm(6.4), height=Cm(2)
        )
        tf_header = header.text_frame
        text_frame_paragraph(tf_header, f"Señor(a) {client['name']}.", 14, True)
        text_frame_paragraph(tf_header, client["company"], 14, True)

    def add_commercial_policy_slide(self):
        self.prs.slides[self.ref_q].shapes.add_picture(
            "./images/condiciones.jpeg",
            left=Cm(1),
            top=Cm(4.5),
            width=Cm(17.27),
            height=Cm(16.66),
        )

    def save(self, path: str):
        self.prs.save(path)
