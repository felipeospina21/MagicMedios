import logging
import re
import time

import requests
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.wait import WebDriverWait

from utils import text_frame_paragraph

logging.basicConfig(
    level=logging.ERROR,
    filename="app.log",
    filemode="w",
    format="%(asctime)s %(name)-12s %(levelname)-8s %(message)s",
    datefmt="%m-%d %H:%M",
)


class Get_Data:
    def __init__(self, supplier, prs, references, measures):
        self.supplier = supplier
        self.prs = prs
        self.references = references

        self.lf_1 = Cm(measures["lf_1"])
        self.lf_2 = Cm(measures["lf_2"])
        self.lf_3 = Cm(measures["lf_3"])
        self.lf_6 = Cm(measures["lf_6"])

        self.t_1 = measures["t_1"]
        self.t_2 = measures["t_2"]
        self.t_3 = measures["t_3"]
        self.t_4 = measures["t_4"]
        self.t_5 = measures["t_5"]
        self.t_6 = measures["t_6"]

        self.w_1 = Cm(measures["w_1"])
        self.w_2 = Cm(measures["w_2"])
        self.w_3 = Cm(measures["w_3"])

        self.h_1 = Cm(measures["h_1"])
        self.h_2 = Cm(measures["h_2"])
        self.h_3 = Cm(measures["h_3"])
        self.h_4 = Cm(measures["h_4"])
        self.h_5 = Cm(measures["h_5"])

        self.cell_font = Pt(measures["cell_font"])
        self.cell_font_2 = Pt(measures["cell_font_2"])

    def error_logging(self):
        logging.error("Error", exc_info=True)
        # exit()

    def execute_driver(self, url):
        service = Service()
        options = webdriver.ChromeOptions()
        options.add_experimental_option("excludeSwitches", ["enable-logging"])

        if self.supplier == "promo_op" or self.supplier == "cat_promo":
            # options.page_load_strategy = 'eager'
            self.driver = webdriver.Chrome(service=service, options=options)
        elif self.supplier == "nw_promo":
            # options.headless = True
            self.driver = webdriver.Chrome(service=service, options=options)
            self.driver.set_page_load_timeout(200)
        else:
            self.driver = webdriver.Chrome(service=service, options=options)
            self.driver.set_page_load_timeout(20)

        try:
            self.driver.get(url)
        except Exception as e:
            self.error_logging()
            raise Exception(e)

        if self.supplier == "cat_promo":
            time.sleep(10)

        time.sleep(5)

    def previous_page(self):
        try:
            self.driver.execute_script("window.history.go(-1)")
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def stop_loading(self):
        try:
            self.driver.execute_script("window.stop();")
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def check_pop_up(self):
        if self.supplier == "nw_promo":
            try:
                self.driver.execute_script(
                    "document.getElementsByClassName('fancybox-overlay fancybox-overlay-fixed labpopup')[0].style.display = 'none';"
                )
            except Exception as e:
                self.error_logging()
                raise Exception(e)

    def get_element_with_xpath(self, xpath):
        try:
            if self.supplier == "cat_promo":
                element = self.driver.find_element(By.XPATH, xpath)
            else:
                element = WebDriverWait(self.driver, 30).until(
                    EC.presence_of_element_located((By.XPATH, xpath))
                )

            return element

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_elements_with_xpath(self, xpath):
        try:
            elements_list = self.driver.find_elements(By.XPATH, xpath)

            return elements_list

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_elements_len_with_xpath(self, xpath):
        try:
            time.sleep(5)
            elements = self.driver.find_elements(By.XPATH, xpath)
            return len(elements)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_element_attribute(self, element, attribute):
        try:
            return element.get_attribute(attribute)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_element_css_property(self, element, css_property):
        try:
            return element.value_of_css_property(css_property)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def fill_stock_table(self, table, color, stock, row_index):
        try:
            c1 = table.cell(row_index, 0)
            c1.text = color
            c1.text_frame.paragraphs[0].font.size = self.cell_font
            c2 = table.cell(row_index, 1)
            c2.text = stock
            c2.text_frame.paragraphs[0].font.size = self.cell_font
            table.rows[row_index].height = Cm(0.5)
            # Cell Color
            cell1 = table.cell(row_index, 0)
            cell2 = table.cell(row_index, 1)
            cell1.fill.solid()
            cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell2.fill.solid()
            cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def send_keys(self, element, text):
        try:
            element.clear()
            element.send_keys(text)
            element.send_keys(Keys.RETURN)
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def accept_alert_popup(self):
        try:
            alert = self.driver.switch_to.alert
            alert.accept()
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def click_first_result(self, first_result_xpath):
        try:
            result = WebDriverWait(self.driver, 15).until(
                EC.element_to_be_clickable((By.XPATH, first_result_xpath))
            )
            time.sleep(1)
            result.click()
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_original_ref_list_idx(self, ref: str) -> int:
        if self.supplier == "cat_promo":
            return self.references.index("CP" + ref)
        elif self.supplier == "mp_promo":
            return self.references.index("MP" + ref)
        elif self.supplier == "promo_op":
            return self.references.index("PO" + ref)
        elif self.supplier == "nw_promo":
            return self.references.index(ref)
        elif self.supplier == "cdo_promo":
            return self.references.index("CD" + ref)
        else:
            raise Exception("not supported supplier")

    def get_title_and_subtitle(self, header_xpath, title_index, subtitle_index):
        try:
            header = self.driver.find_element(By.XPATH, header_xpath)
            header_text = header.text.split("\n")
            title = header_text[title_index]
            subtitle = header_text[subtitle_index]
            return (title, subtitle)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_title_with_xpath(self, title_xpath):
        try:
            title = WebDriverWait(self.driver, 10).until(
                EC.presence_of_element_located((By.XPATH, title_xpath))
            )
            title = self.driver.find_element(By.XPATH, title_xpath)
            return title.text

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_subtitle_with_xpath(self, sub_title_xpath):
        def get_subtitle_promo_op(subtitle_result):
            split_text = subtitle_result.text.split("\n")

            result_idx = 0
            for txt in split_text:
                if re.search("^Desc", txt):
                    result_idx = split_text.index(txt)

            result = split_text[result_idx].split("Descripción: ")
            return result[1]

        try:
            subtitle = WebDriverWait(self.driver, 10).until(
                EC.presence_of_element_located((By.XPATH, sub_title_xpath))
            )
            subtitle = self.driver.find_element(By.XPATH, sub_title_xpath)
            if self.supplier == "promo_op":
                subtitle_text = get_subtitle_promo_op(subtitle)
                return subtitle_text

            return subtitle.text

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_description(self, desc_xpath):
        try:
            desc = self.driver.find_elements(By.XPATH, desc_xpath)
            return desc

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def get_package_info(self, ref):
        try:
            table = "//table[@class='table-list']"
            unit_col_1 = self.driver.find_element(
                By.XPATH, f"{table}/tbody[1]/tr[1]/td[1]"
            )
            unit_col_2 = self.driver.find_element(
                By.XPATH, f"{table}/tbody[1]/tr[1]/td[2]"
            )
            package_col_1 = self.driver.find_element(
                By.XPATH, f"{table}/tbody[1]/tr[2]/td[1]"
            )
            package_col_2 = self.driver.find_element(
                By.XPATH, f"{table}/tbody[1]/tr[2]/td[2]"
            )
            table_texts_list = [
                unit_col_1.text,
                unit_col_2.text,
                package_col_1.text,
                package_col_2.text,
            ]

            return table_texts_list

        except Exception as e:
            print(f"Error de tipo {e.__class__}")
            print(f"No se pudo obtener la info de empaque de la ref {ref}")

    def get_img(self, img_xpath):
        try:
            if self.supplier == "promo_op":
                time.sleep(5)

            img = self.driver.find_element(By.XPATH, img_xpath)
            img_src = img.get_attribute("src")

            return img_src

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_quantity_table(self, idx):
        ROWS = 3
        COLS = 4

        def createHeader(cell, text):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(26, 152, 139)

        def createRowCell(cell, text):
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        try:
            if idx > 0:
                top = Cm(self.t_5 - 1)
            else:
                top = Cm(self.t_5)

            table = (
                self.prs.slides[idx]
                .shapes.add_table(ROWS, COLS, self.lf_6, top, self.w_1, self.h_2)
                .table
            )

            c1 = table.cell(0, 0)
            c2 = table.cell(0, 1)
            c3 = table.cell(0, 2)
            c4 = table.cell(0, 3)

            createHeader(c1, "CANTIDAD")
            createHeader(c2, "TÉCNICA DE MARCACIÓN")
            createHeader(c3, "DETALLE")
            createHeader(c4, "VALOR UNITARIO ANTES DE IVA")

            for i in range(0, COLS):
                table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            table.rows[0].height = Cm(0.5)
            table.first_row = True
            table.horz_banding = False
            for i in range(1, 3):
                table.rows[i].height = Cm(0.5)
                # Cell Color
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                cell3 = table.cell(i, 2)
                cell4 = table.cell(i, 3)

                createRowCell(cell1, "(Und)")
                createRowCell(cell2, "")
                createRowCell(cell3, "")
                createRowCell(cell4, "$")

            table.columns[0].width = Cm(3)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_title(self, title_text, idx, count, ref):
        try:
            title = title_text
            if idx > 0:
                top = Cm(self.t_1 - 1)
            else:
                top = Cm(self.t_1)

            titulo = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_1, height=self.h_1
            )
            tf_titulo = titulo.text_frame
            text_frame_paragraph(tf_titulo, f"{count}. {title} {ref}", 12, True)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_subtitle(self, subtitle_text, idx):
        try:
            subtitle = subtitle_text
            if idx > 0:
                top = Cm(self.t_2 - 1)
            else:
                top = Cm(self.t_2)

            sub_titulo = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_1, height=self.h_2
            )
            tf_sub_titulo = sub_titulo.text_frame
            tf_sub_titulo.word_wrap = True
            text_frame_paragraph(tf_sub_titulo, subtitle, 11)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_description(self, desc_list, idx):
        try:
            if idx > 0:
                top = Cm(self.t_3 - 1)
            else:
                top = Cm(self.t_3)

            description = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_1, height=self.h_3
            )
            tf_desc = description.text_frame
            tf_desc.word_wrap = True
            for element in desc_list:
                text_frame_paragraph(tf_desc, element.text, 11)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    # TODO: Replace this function for create_description.
    # Identical implementation, here desc_list is a string[]
    def create_desc(self, desc_list, idx):
        try:
            if idx > 0:
                top = Cm(self.t_3 - 1)
            else:
                top = Cm(self.t_3)

            description = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_1, height=self.h_3
            )
            tf_desc = description.text_frame
            tf_desc.word_wrap = True
            for element in desc_list:
                text_frame_paragraph(tf_desc, element, 11)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_description_promo_op(self, desc_list, idx):
        try:
            if idx > 0:
                top = Cm(self.t_3 - 1)
            else:
                top = Cm(self.t_3)

            description_1 = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_3, height=self.h_3
            )
            description_2 = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_2, top=top, width=self.w_3, height=self.h_3
            )
            tf_desc_1 = description_1.text_frame
            tf_desc_1.word_wrap = True
            tf_desc_2 = description_2.text_frame
            tf_desc_2.word_wrap = True

            list_len = len(desc_list)
            limit = int(list_len / 2)

            for i in range(0, limit - 1):
                text_frame_paragraph(tf_desc_1, desc_list[i].text, 11)
            for i in range(limit, list_len - 1):
                text_frame_paragraph(tf_desc_2, desc_list[i].text, 11)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_printing_info(self, printing_methods_list, idx):
        try:
            if idx > 0:
                top = Cm(self.t_4 - 1)
            else:
                top = Cm(self.t_4)

            p1 = self.prs.slides[idx].shapes.add_textbox(
                left=self.lf_1, top=top, width=self.w_1, height=self.h_2
            )
            tf_p1 = p1.text_frame

            for element in printing_methods_list:
                text_frame_paragraph(tf_p1, element, 11)

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_inventory_table(self, q_colores, xpath_tabla_colores, idx):
        try:
            cols = 2
            rows = q_colores
            if idx > 0:
                top = Cm(self.t_6 - 1)
            else:
                top = Cm(self.t_6)

            table = (
                self.prs.slides[idx]
                .shapes.add_table(rows + 1, cols, self.lf_1, top, self.w_2, self.h_4)
                .table
            )

            # Table Header
            h1 = table.cell(0, 0)
            h2 = table.cell(0, 1)
            h1.text = "Color"
            h2.text = "Inventario"
            h1.text_frame.paragraphs[0].font.size = self.cell_font
            h2.text_frame.paragraphs[0].font.size = self.cell_font
            table.rows[0].height = Cm(0.5)
            table.first_row = False
            table.horz_banding = False

            for i in range(1, q_colores + 1):
                if self.supplier == "cat_promo":
                    color_xpath = f"tbody[1]/tr[not(@class='hideInfo')][{i+2}]/td[1]"
                    inv_color_xpath = (
                        f"tbody[1]/tr[not(@class='hideInfo')][{i+2}]/td[4]"
                    )

                elif self.supplier == "mp_promo":
                    color_xpath = f"tr[{i}]/td[3]"
                    inv_color_xpath = f"tr[{i}]/td[6]"

                elif self.supplier == "nw_promo":
                    color_xpath = f"tr[{i}]/td[1]"
                    inv_color_xpath = f"tr[{i}]/td[5]"

                else:
                    raise Exception("Not supported method for supplier")

                color = self.driver.find_element(
                    By.XPATH, f"{xpath_tabla_colores}/{color_xpath}"
                ).text
                inv_color = self.driver.find_element(
                    By.XPATH, f"{xpath_tabla_colores}/{inv_color_xpath}"
                ).text
                c1 = table.cell(i, 0)
                c1.text = color
                c1.text_frame.paragraphs[0].font.size = self.cell_font
                c2 = table.cell(i, 1)
                c2.text = inv_color
                c2.text_frame.paragraphs[0].font.size = self.cell_font
                table.rows[i].height = Cm(0.5)
                # Cell Color
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                cell1.fill.solid()
                cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell2.fill.solid()
                cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)

            table.columns[0].width = Cm(3.8)
            table.columns[1].width = Cm(2.2)
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_stock_table_api(self, q_colores, colors_list, idx):
        try:
            cols = 2
            rows = q_colores
            if idx > 0:
                top = Cm(self.t_6 - 1)
            else:
                top = Cm(self.t_6)
            table = (
                self.prs.slides[idx]
                .shapes.add_table(rows + 1, cols, self.lf_1, top, self.w_2, self.h_4)
                .table
            )

            # Table Header
            h1 = table.cell(0, 0)
            h2 = table.cell(0, 1)
            h1.text = "Color"
            h2.text = "Inventario"
            h1.text_frame.paragraphs[0].font.size = self.cell_font
            h2.text_frame.paragraphs[0].font.size = self.cell_font
            table.rows[0].height = Cm(0.5)
            table.first_row = False
            table.horz_banding = False

            for i in range(1, q_colores + 1):
                color = colors_list[i - 1]["color"]
                stock = str(colors_list[i - 1]["stock_available"])
                c1 = table.cell(i, 0)
                c1.text = color
                c1.text_frame.paragraphs[0].font.size = self.cell_font
                c2 = table.cell(i, 1)
                c2.text = stock
                c2.text_frame.paragraphs[0].font.size = self.cell_font
                table.rows[i].height = Cm(0.5)
                # Cell Color
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                cell1.fill.solid()
                cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell2.fill.solid()
                cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)

            table.columns[0].width = Cm(3.8)
            table.columns[1].width = Cm(2.2)
        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_stock_table(self, colors_q, idx):
        try:
            cols = 2
            rows = colors_q
            if idx > 0:
                top = Cm(self.t_6 - 1)
            else:
                top = Cm(self.t_6)
            table = (
                self.prs.slides[idx]
                .shapes.add_table(rows, cols, self.lf_1, top, self.w_2, self.h_4)
                .table
            )

            # Table Header
            h1 = table.cell(0, 0)
            h2 = table.cell(0, 1)
            h1.text = "Color"
            h2.text = "Inventario"
            h1.text_frame.paragraphs[0].font.size = self.cell_font
            h2.text_frame.paragraphs[0].font.size = self.cell_font
            table.rows[0].height = Cm(0.5)
            table.first_row = False
            table.horz_banding = False

            table.columns[0].width = Cm(3.8)
            table.columns[1].width = Cm(2.2)

            return table

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def create_img(self, img_src, idx, img_height, ref):
        response = requests.get(img_src)
        try:
            if idx > 0:
                top = Cm(self.t_6 - 1)
            else:
                top = Cm(self.t_6)

            if response.status_code == 200:
                file = open("./images/sample_image.jpg", "wb")
                file.write(response.content)
                file.close()

                if self.supplier == "cdo_promo":
                    self.prs.slides[idx].shapes.add_picture(
                        "./images/sample_image.jpg", left=self.lf_2, top=top
                    )

                else:
                    self.prs.slides[idx].shapes.add_picture(
                        "./images/sample_image.jpg",
                        left=self.lf_2,
                        top=top,
                        # width=Cm(img_width),
                        height=Cm(img_height),
                    )

            else:
                print(
                    f"Error al descargar imagen de la ref {ref}, status code({response.status_code})"
                )

        except Exception as e:
            self.error_logging()
            raise Exception(e)

    def close_driver(self):
        self.driver.close()
