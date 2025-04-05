import asyncio
import json
import re

import requests
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from playwright.async_api import Locator, Page


def create_supplier_ref_list(
    ref, suppliers_dict: dict[str, list[str]]
) -> dict[str, list[str]]:
    if re.search("^CP|^cp]", ref):
        split_ref = ref.split("CP", 1)
        suppliers_dict["cat_promo"].append(split_ref[1])
    elif re.search("^MP|^mp]", ref):
        split_ref = ref.split("MP", 1)
        suppliers_dict["mp_promo"].append(split_ref[1])
    elif re.search("^PO|^po]", ref):
        split_ref = ref.split("PO", 1)
        suppliers_dict["promo_op"].append(split_ref[1])
    elif re.search("^CD|^cd", ref):
        split_ref = ref.split("CD", 1)
        suppliers_dict["cdo_promo"].append(split_ref[1])
    elif re.search("^NW|^nw", ref):
        suppliers_dict["nw_promo"].append(ref)
    else:
        print(
            f"\nNo se pudo identificar la referencia {ref}, favor validar el prefijo ingresado"
        )

    return suppliers_dict


def text_frame_paragraph(text_frame, text, font_size, bold=False, centered=False):
    tf = text_frame.add_paragraph()
    tf.text = text
    tf.font.size = Pt(font_size)
    tf.font.bold = bold
    tf.space_before = Cm(0)
    if centered:
        tf.alignment = PP_ALIGN.CENTER


def get_api_data(url):
    response = requests.get(url)
    content = response.content
    return json.loads(content)


measures: dict[str, float] = {
    "lf_1": 0.8,
    "lf_2": 8.5,
    "lf_3": 3,
    "lf_4": 12,
    "lf_5": 1,
    "lf_6": 1.5,
    "t_0": -0.5,
    "t_1": 3.5,
    "t_2": 4,
    "t_3": 5.5,
    "t_4": 8.5,
    "t_5": 11.5,
    "t_6": 14.5,
    "t_7": 2.5,
    "w_1": 17.4,
    "w_2": 6,
    "w_3": 8,
    "w_4": 6.6,
    "w_5": 6.4,
    "h_1": 1,
    "h_2": 2,
    "h_3": 5,
    "h_4": 8.5,
    "h_5": 8,
    "h_6": 6,
    "cell_font": 7,
    "cell_font_2": 11,
}


async def wait_for_selector_with_retry(
    page: Page, selector: str, timeout: int = 5000, retries: int = 3, delay: int = 2
) -> bool:
    """Retries waiting for a selector multiple times before failing."""
    for attempt in range(1, retries + 1):
        try:
            element = page.locator(selector)
            await element.wait_for(timeout=timeout)
            print(f"wating for selector {selector}, {attempt}/{retries}")
            return True
        except Exception as e:
            if attempt < retries - 1:
                await asyncio.sleep(delay)
            else:
                return False
    return False


async def get_selector_with_retry(
    page: Page, selector: str, timeout: int = 5000, retries: int = 3, delay: int = 2
) -> Locator | None:
    """Retries waiting for a selector multiple times before failing."""
    for attempt in range(1, retries + 1):
        try:
            print(f"wating for selector {selector}, {attempt}/{retries}")
            element = page.locator(selector)
            await element.wait_for(timeout=timeout)
            return element
        except Exception:
            if attempt < retries - 1:
                await asyncio.sleep(delay)
            else:
                return None


async def get_all_selectors_with_retry(
    page: Page, selector: str, timeout: int = 1000, retries: int = 3, delay: int = 0
) -> list[Locator] | None:
    """Retries waiting for a selector multiple times before failing."""
    for attempt in range(1, retries + 1):
        try:
            print(f"wating for selectors {selector}, {attempt}/{retries}")
            elements = await page.locator(selector).all()
            for element in elements:
                await element.wait_for(timeout=timeout)
            return elements
        except Exception:
            if attempt < retries - 1:
                await asyncio.sleep(delay)
            else:
                return None


def create_inventory_table(colors, xpath, idx):
    try:
        cols = 2
        rows = colors
        if idx > 0:
            top = Cm(self.t_6 - 1)
        else:
            top = Cm(self.t_6)

        table = (
            self.prs.slides[idx]
            .shapes.add_table(rows + 1, cols, self.lf_1, top, self.w_2, self.h_4)
            .table
        )

        # Table Header
        h1 = table.cell(0, 0)
        h2 = table.cell(0, 1)
        h1.text = "Color"
        h2.text = "Inventario"
        h1.text_frame.paragraphs[0].font.size = self.cell_font
        h2.text_frame.paragraphs[0].font.size = self.cell_font
        table.rows[0].height = Cm(0.5)
        table.first_row = False
        table.horz_banding = False

        for i in range(1, colors + 1):
            if self.supplier == "cat_promo":
                color_xpath = f"tbody[1]/tr[not(@class='hideInfo')][{i+2}]/td[1]"
                inv_color_xpath = f"tbody[1]/tr[not(@class='hideInfo')][{i+2}]/td[4]"

            elif self.supplier == "mp_promo":
                color_xpath = f"tr[{i}]/td[3]"
                inv_color_xpath = f"tr[{i}]/td[6]"

            elif self.supplier == "nw_promo":
                color_xpath = f"tr[{i}]/td[1]"
                inv_color_xpath = f"tr[{i}]/td[5]"

            else:
                raise Exception("Not supported method for supplier")

            color = self.driver.find_element(By.XPATH, f"{xpath}/{color_xpath}").text
            inv_color = self.driver.find_element(
                By.XPATH, f"{xpath}/{inv_color_xpath}"
            ).text
            c1 = table.cell(i, 0)
            c1.text = color
            c1.text_frame.paragraphs[0].font.size = self.cell_font
            c2 = table.cell(i, 1)
            c2.text = inv_color
            c2.text_frame.paragraphs[0].font.size = self.cell_font
            table.rows[i].height = Cm(0.5)
            # Cell Color
            cell1 = table.cell(i, 0)
            cell2 = table.cell(i, 1)
            cell1.fill.solid()
            cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell2.fill.solid()
            cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)

        table.columns[0].width = Cm(3.8)
        table.columns[1].width = Cm(2.2)
    except Exception as e:
        self.error_logging(e)
        raise SystemExit("Error: ", e)
