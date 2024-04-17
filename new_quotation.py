from utils import text_frame_paragraph, create_supplier_ref_list
from catalogospromo import get_cat_promo_data
from mppromos import get_mp_promo_data
from nwpromo import get_nw_promo_data
from promoop import get_promo_op__data
from cdopromo import get_cdo_promo_data
from pptx.util import Cm
import os


class New_Quotation:
    def __init__(self, prs, title_slide_layout, username) -> None:
        if username == "felipe" or username == "felip":
            self.path = os.environ.get("FILE_PATH")
            self.path = "."
            self.consecutive_path = f"{self.path}/data/consecutivo.txt"
        else:
            self.path = os.environ.get("COTIZACIONES_PATH")
            self.consecutive_path = f"{self.path}/Z consecutivo.txt"

        self.user_path = f"{self.path}/data"
        self.footer_path = f"{self.path}/data/data.txt"
        self.users = self.get_users()
        self.client = input("Ingrese nombre cliente: ").title()
        self.company = input("Ingrese nombre empresa: ").upper()
        self.user = input(
            f"Ingrese nombre asesor ({', '.join(self.users)}): ").lower()
        self.reference = (
            input("Ingrese referencias a consultar (separadas por coma): ")
            .upper()
            .split(",")
        )

        self.consecutive = self.get_consecutive()
        self.rep_name, self.phone, self.email = self.get_contact_info()
        self.address, self.web = self.get_footer()
        self.prs = prs
        self.title_slide_layout = title_slide_layout

        self.strip_reference = [ref.strip() for ref in self.reference]
        self.ref_q = len(self.reference)

        if username == "felipe" or username == "felip":
            self.save_path = f"{self.path}/cotizaciones/cotización_{self.company}.pptm"
        else:
            self.save_path = f"{self.path}/Cotización N°{self.consecutive} - {self.company} - Magic Medios SAS.pptm"

    def get_users(self):
        users = []
        exclude = ["consecutivo", "data"]
        for x in os.listdir(self.user_path):
            if x.endswith(".txt"):
                user = x.split(".txt")[0]
                if user not in exclude:
                    users.append(user)

        return users

    def get_consecutive(self):
        file = open(self.consecutive_path, "r")
        consecutive = file.readline().strip()
        file.close()
        return consecutive

    def get_contact_info(self):
        file = open(f"{self.user_path}/{self.user}.txt", "r")
        rep_name = file.readline()
        phone = file.readline()
        email = file.readline()
        file.close()

        return rep_name, phone, email

    def get_footer(self):
        file = open(self.footer_path, "r")
        address = file.readline()
        web = file.readline()
        file.close()
        return address, web

    def set_slides(self):
        for i in range(0, self.ref_q + 1):
            self.prs.slides.add_slide(self.title_slide_layout)
            self.prs.slides[i].shapes.add_picture(
                "./images/logo.jpeg",
                left=Cm(1),
                top=Cm(0.5),
                width=Cm(8.9),
                height=Cm(1.7),
            )
            footer = self.prs.slides[i].shapes.add_textbox(
                left=Cm(0.5), top=Cm(22.8), width=Cm(18), height=Cm(1)
            )
            tf_footer = footer.text_frame
            text_frame_paragraph(
                tf_footer,
                f"{self.address} {self.phone} {self.email} {self.web}",
                7,
                False,
                True,
            )

    def add_header(self, today):
        txBox = self.prs.slides[0].shapes.add_textbox(
            left=Cm(12), top=Cm(-0.5), width=Cm(6.6), height=Cm(6)
        )
        tf = txBox.text_frame
        text_frame_paragraph(
            tf, f"{today.day} {today.strftime('%B')} de {today.year}", 14
        )
        text_frame_paragraph(tf, f"Cot N°{self.consecutive}", 14)
        text_frame_paragraph(tf, "", 11)
        text_frame_paragraph(tf, "Asesor Comercial", 11)
        text_frame_paragraph(
            tf, f"{self.rep_name} {self.phone} {self.email}", 11)

    def add_client_name(self):
        header = self.prs.slides[0].shapes.add_textbox(
            left=Cm(1), top=Cm(1.8), width=Cm(6.4), height=Cm(2)
        )
        tf_header = header.text_frame
        text_frame_paragraph(tf_header, f"Señor(a) {self.client}.", 14, True)
        text_frame_paragraph(tf_header, self.company, 14, True)

    def create_suppliers_ref_list(self, suppliers):
        for ref in self.strip_reference:
            suppliers_list = create_supplier_ref_list(ref, suppliers)

        return suppliers_list

    def create_commercial_policy(self):
        self.prs.slides[self.ref_q].shapes.add_picture(
            "./images/condiciones.jpeg",
            left=Cm(1),
            top=Cm(4.5),
            width=Cm(17.27),
            height=Cm(16.66),
        )

    def create_new_consecutive(self):
        new_consecutive = int(self.consecutive) + 1
        file = open(self.consecutive_path, "w")
        file.write(f"{new_consecutive}\n")
        file.close()

    def save(self):
        self.prs.save(self.save_path)

    def scrapp_data(self, suppliers_list, suppliers_dict):
        if len(suppliers_list["cat_promo"]) != 0:
            get_cat_promo_data(suppliers_dict, self.prs, self.strip_reference)
        if len(suppliers_list["mp_promo"]) != 0:
            get_mp_promo_data(suppliers_dict, self.prs, self.strip_reference)
        if len(suppliers_list["promo_op"]) != 0:
            get_promo_op__data(suppliers_dict, self.prs, self.strip_reference)
        if len(suppliers_list["nw_promo"]) != 0:
            get_nw_promo_data(suppliers_dict, self.prs, self.strip_reference)
        if len(suppliers_list["cdo_promo"]) != 0:
            get_cdo_promo_data(suppliers_dict, self.prs, self.strip_reference)
