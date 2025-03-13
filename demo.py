import asyncio
import requests
from io import BytesIO
from PIL import Image
from demo_create_ppt import Get_Data
from playwright.async_api import async_playwright, Browser, Page
from pptx import Presentation
from pptx.util import Inches
from asyncio import Semaphore
from typing import Optional, Sequence, Tuple, List
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


type Data = Tuple[str, str, list[str], Optional[BytesIO]]

measures = {
    "lf_1": 0.8,
    "lf_2": 8.5,
    "lf_3": 3,
    "lf_4": 12,
    "lf_5": 1,
    "lf_6": 1.5,
    "t_0": -0.5,
    "t_1": 3.5,
    "t_2": 4,
    "t_3": 5.5,
    "t_4": 8.5,
    "t_5": 11.5,
    "t_6": 14.5,
    "t_7": 2.5,
    "w_1": 17.4,
    "w_2": 6,
    "w_3": 8,
    "w_4": 6.6,
    "w_5": 6.4,
    "h_1": 1,
    "h_2": 2,
    "h_3": 5,
    "h_4": 8.5,
    "h_5": 8,
    "h_6": 6,
    "cell_font": 7,
    "cell_font_2": 11,
}


async def scrape_product_with_limit(
    sem: Semaphore, browser: Browser, product_code: str
) -> Data:
    """Limits concurrent product scrapes."""
    async with sem:
        return await scrape_product(browser, product_code)


async def scrape_product(browser: Browser, product_code: str) -> Data:
    """Scrapes product name and image from the website."""
    context = await browser.new_context()
    page = await context.new_page()

    await page.goto("https://www.catalogospromocionales.com/")
    await page.wait_for_load_state("domcontentloaded")  # Waits for network to be idle
    print(f"Processing: {product_code}")

    found: bool = await search_product(page, product_code, delay=0)
    if not found:
        await context.close()
        return product_code, "Not Found", [], None

    await page.click(".img-producto")

    found = await wait_for_selector_with_retry(page, "#img_01", timeout=5000)
    if not found:
        await context.close()
        return product_code, "Not Found", [], None

    product_name: str = await page.locator(".hola").inner_text()
    product_image_url: Optional[str] = await page.locator(
        "#img_01"
    ).first.get_attribute("src")
    info_text_arr = product_name.split("\n\n")
    title = info_text_arr[0]
    desc_list = info_text_arr[2:]

    if not product_image_url:
        await context.close()
        return product_code, title, desc_list, None

    response = requests.get(product_image_url)
    image_data = BytesIO(response.content)

    await context.close()
    return product_code, title, desc_list, image_data


async def search_product(
    page: Page, product_code: str, retries: int = 3, delay: int = 2
) -> bool:
    """Attempts to search for a product, retrying if necessary."""
    for attempt in range(retries):
        input: bool = await wait_for_selector_with_retry(
            page, "#productos", retries=3, delay=0
        )
        if input:
            await asyncio.sleep(2)
            await page.locator("#productos").fill(product_code)
            await page.locator("#productos").press("Enter")

        found: bool = await wait_for_selector_with_retry(
            page, ".img-producto", retries=3, delay=0
        )
        if found:
            return True

        await asyncio.sleep(delay)

    return False


async def wait_for_selector_with_retry(
    page: Page, selector: str, timeout: int = 5000, retries: int = 3, delay: int = 2
) -> bool:
    """Retries waiting for a selector multiple times before failing."""
    for attempt in range(retries):
        try:
            await page.wait_for_selector(selector, timeout=timeout)
            print(f"wating for selector {selector}, {attempt}/{retries}")
            return True
        except Exception:
            if attempt < retries - 1:
                await asyncio.sleep(delay)
            else:
                return False
    return False


def text_frame_paragraph(text_frame, text, font_size, bold=False, centered=False):
    tf = text_frame.add_paragraph()
    tf.text = text
    tf.font.size = Pt(font_size)
    tf.font.bold = bold
    tf.space_before = Cm(0)
    if centered:
        tf.alignment = PP_ALIGN.CENTER


async def create_ppt(products: Sequence[Data], prs, original_refs) -> None:
    """Creates a PowerPoint file with product details."""

    data = Get_Data("cat_promo", prs, original_refs, measures)
    for product_code, product_name, desc, image_data in products:
        idx = data.get_original_ref_list_idx(product_code)
        count = idx + 1
        data.create_quantity_table(idx)
        data.create_title(product_name, idx, count, product_code)
        data.create_desc(desc, idx)
        data.create_img(image_data, idx)
        # slide = prs.slides.add_slide(prs.slide_layouts[5])
        #
        # title = slide.shapes.title
        # title.text = f"{product_code} - {product_name}"
        #
        # description = slide.shapes.add_textbox(
        #     left=measures["lf_1"],
        #     top=measures["t_1"],
        #     width=measures["w_1"],
        #     height=measures["h_3"],
        # )
        # tf_desc = description.text_frame
        # tf_desc.word_wrap = True
        # for element in desc:
        #     text_frame_paragraph(tf_desc, element, 11)
        #
        # if image_data:
        #     left, top, height = Inches(1), Inches(2), Inches(3)
        #     image = Image.open(image_data)
        #     image_bytes = BytesIO()
        #     image.save(image_bytes, format="PNG")
        #     prs.slides[idx].shapes.add_picture(image_bytes, left, top, height=height)

    prs.save("Product_Presentation.pptx")
    print("PowerPoint created: Product_Presentation.pptx")


async def process_products(
    product_codes: List[str],
    prs,
    original_refs: List[str],
    max_concurrent: int = 5,
) -> None:
    """Processes multiple products concurrently with a limit."""
    sem = Semaphore(max_concurrent)

    async with async_playwright() as p:
        browser: Browser = await p.chromium.launch(headless=False)

        tasks = [
            scrape_product_with_limit(sem, browser, code) for code in product_codes
        ]
        results: List[Data] = await asyncio.gather(*tasks)

        await browser.close()

    product_data: List[Data] = [
        (code, name, desc, img) for code, name, desc, img in results if img
    ]

    await create_ppt(product_data, prs, original_refs)


async def get_selector(selector: str, product_code: str, page, context):
    found = await wait_for_selector_with_retry(page, selector, timeout=5000)
    if not found:
        await context.close()
        return product_code, "Not Found", None


async def main(refs: List[str], prs, original_refs: List[str]):
    product_codes: List[str] = [
        "VA-1153",
        "VA-509",
        "SO-65",
        "SO-75",
        "SO-47",
        "SO-44",
        "va-1148",
        "va-1147",
        "va-1149",
        "so-45",
        "so-41",
        "so-43",
        "so-50",
        "so-56",
        "so-58",
    ]  # Add more product codes

    await process_products(refs, prs, original_refs, max_concurrent=5)


# asyncio.run(main())
