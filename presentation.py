from datetime import datetime
from io import BytesIO

from PIL import Image
from pptx import Presentation as PPTX
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from constants import measures
from entities.entities import (
    Client,
    Color_Inventory,
    Contact,
    ProductData,
    Representative,
)


class Presentation:
    prs = PPTX("./plantillas/cotizacion.pptm")
    title_slide_layout = prs.slide_layouts[6]

    def __init__(self, references_quantity: int) -> None:
        self.ref_q = references_quantity

    def text_frame_paragraph(
        self, text_frame, text, font_size, bold=False, centered=False
    ):
        tf = text_frame.add_paragraph()
        tf.text = text
        tf.font.size = Pt(font_size)
        tf.font.bold = bold
        tf.space_before = Cm(0)
        if centered:
            tf.alignment = PP_ALIGN.CENTER

    def set_slides(self, contact: Contact):
        for i in range(0, self.ref_q + 1):
            self.prs.slides.add_slide(self.title_slide_layout)
            self.prs.slides[i].shapes.add_picture(
                "./images/logo.jpeg",
                left=Cm(1),
                top=Cm(0.5),
                width=Cm(8.9),
                height=Cm(1.7),
            )
            footer_textbox = self.prs.slides[i].shapes.add_textbox(
                left=Cm(0.5), top=Cm(22.8), width=Cm(18), height=Cm(1)
            )
            tf_footer = footer_textbox.text_frame
            self.text_frame_paragraph(
                tf_footer,
                f"{contact['address']} {contact['phone']} {contact['email']} {contact['web']}",
                7,
                False,
                True,
            )

    def add_header(self, representative: Representative, consecutive: int):
        today = datetime.now()
        txBox = self.prs.slides[0].shapes.add_textbox(
            left=Cm(12), top=Cm(-0.5), width=Cm(6.6), height=Cm(6)
        )
        tf = txBox.text_frame
        self.text_frame_paragraph(
            tf, f"{today.day} {today.strftime('%B')} de {today.year}", 14
        )
        self.text_frame_paragraph(tf, f"Cot N°{consecutive}", 14)
        self.text_frame_paragraph(tf, "", 11)
        self.text_frame_paragraph(tf, "Asesor Comercial", 11)
        self.text_frame_paragraph(
            tf,
            f"{representative['name']} {representative['phone']} {representative['email']}",
            11,
        )

    def add_client_name(self, client: Client):
        header = self.prs.slides[0].shapes.add_textbox(
            left=Cm(1), top=Cm(1.8), width=Cm(6.4), height=Cm(2)
        )
        tf_header = header.text_frame
        self.text_frame_paragraph(tf_header, f"Señor(a) {client['name']}.", 14, True)
        self.text_frame_paragraph(tf_header, client["company"], 14, True)

    def add_commercial_policy_slide(self):
        self.prs.slides[self.ref_q].shapes.add_picture(
            "./images/condiciones.jpeg",
            left=Cm(1),
            top=Cm(4.5),
            width=Cm(17.27),
            height=Cm(16.66),
        )

    def create_title(self, title_text: str, idx: int, ref: str):
        try:
            title = title_text
            if idx > 0:
                top = Cm(measures["t_1"] - 1)
            else:
                top = Cm(measures["t_1"])

            titulo = self.prs.slides[idx].shapes.add_textbox(
                left=Cm(measures["lf_1"]),
                top=top,
                width=Cm(measures["w_1"]),
                height=Cm(measures["h_1"]),
            )
            tf_titulo = titulo.text_frame
            self.text_frame_paragraph(tf_titulo, f"{idx+1}. {title} {ref}", 12, True)

        except Exception as e:
            raise SystemExit("Error: ", e)

    def create_subtitle(self, subtitle_text, idx):
        try:
            subtitle = subtitle_text
            if idx > 0:
                top = Cm(measures["t_2"] - 1)
            else:
                top = Cm(measures["t_2"])

            sub_titulo = self.prs.slides[idx].shapes.add_textbox(
                left=Cm(measures["lf_1"]),
                top=top,
                width=Cm(measures["w_1"]),
                height=Cm(measures["h_2"]),
            )
            tf_sub_titulo = sub_titulo.text_frame
            tf_sub_titulo.word_wrap = True
            self.text_frame_paragraph(tf_sub_titulo, subtitle, 11)

        except Exception as e:
            raise SystemExit("Error: ", e)

    def create_description(self, desc_list, idx):
        try:
            if idx > 0:
                top = Cm(measures["t_3"] - 1)
            else:
                top = Cm(measures["t_3"])

            description = self.prs.slides[idx].shapes.add_textbox(
                left=Cm(measures["lf_1"]),
                top=top,
                width=Cm(measures["w_1"]),
                height=Cm(measures["h_3"]),
            )
            tf_desc = description.text_frame
            tf_desc.word_wrap = True
            for element in desc_list:
                self.text_frame_paragraph(tf_desc, element, 11)

        except Exception as e:
            raise SystemExit("Error: ", e)

    def create_img(self, image_data: BytesIO | None, idx: int):
        if image_data:
            if idx > 0:
                top = Cm(measures["t_6"] - 1)
            else:
                top = Cm(measures["t_6"])
            image = Image.open(image_data)
            image_bytes = BytesIO()
            image.save(image_bytes, format="PNG")
            self.prs.slides[idx].shapes.add_picture(
                image_bytes,
                left=Cm(measures["lf_2"]),
                top=top,
                height=Cm(8),
            )

            # if self.supplier == "cdo_promo":
            #     self.prs.slides[idx].shapes.add_picture(
            #         image_bytes, left=self.lf_2, top=top
            #     )
            #
            # else:
            #     self.prs.slides[idx].shapes.add_picture(
            #         image_bytes,
            #         left=Cm(measures["lf_2"]),
            #         top=top,
            #         height=Cm(8),
            #     )

    def create_quantity_table(self, idx):
        ROWS = 3
        COLS = 4

        def createHeader(cell, text):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(26, 152, 139)

        def createRowCell(cell, text):
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        try:
            if idx > 0:
                top = Cm(measures["t_5"] - 1)
            else:
                top = Cm(measures["t_5"])

            table = (
                self.prs.slides[idx]
                .shapes.add_table(
                    ROWS,
                    COLS,
                    Cm(measures["lf_6"]),
                    top,
                    Cm(measures["w_1"]),
                    Cm(measures["h_2"]),
                )
                .table
            )

            c1 = table.cell(0, 0)
            c2 = table.cell(0, 1)
            c3 = table.cell(0, 2)
            c4 = table.cell(0, 3)

            createHeader(c1, "CANTIDAD")
            createHeader(c2, "TÉCNICA DE MARCACIÓN")
            createHeader(c3, "DETALLE")
            createHeader(c4, "VALOR UNITARIO ANTES DE IVA")

            for i in range(0, COLS):
                table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            table.rows[0].height = Cm(0.5)
            table.first_row = True
            table.horz_banding = False
            for i in range(1, 3):
                table.rows[i].height = Cm(0.5)
                # Cell Color
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                cell3 = table.cell(i, 2)
                cell4 = table.cell(i, 3)

                createRowCell(cell1, "(Und)")
                createRowCell(cell2, "")
                createRowCell(cell3, "")
                createRowCell(cell4, "$")

            table.columns[0].width = Cm(3)

        except Exception as e:
            raise SystemExit("Error: ", e)

    def create_inventory_table(self, inventory: list[Color_Inventory], idx: int):
        try:
            cols = 2
            rows = len(inventory)
            if idx > 0:
                top = Cm(measures["t_6"] - 1)
            else:
                top = Cm(measures["t_6"])

            table = (
                self.prs.slides[idx]
                .shapes.add_table(
                    rows + 1,
                    cols,
                    Cm(measures["lf_1"]),
                    top,
                    Cm(measures["w_2"]),
                    Cm(measures["h_4"]),
                )
                .table
            )

            # Table Header
            h1 = table.cell(0, 0)
            h2 = table.cell(0, 1)
            h1.text = "Color"
            h2.text = "Inventario"
            h1.text_frame.paragraphs[0].font.size = Pt(measures["cell_font"])
            h2.text_frame.paragraphs[0].font.size = Pt(measures["cell_font"])
            table.rows[0].height = Cm(0.5)
            table.first_row = False
            table.horz_banding = False

            for i, element in enumerate(inventory, 1):
                c1 = table.cell(i, 0)
                c1.text = element["color"]
                c1.text_frame.paragraphs[0].font.size = Pt(measures["cell_font"])
                c2 = table.cell(i, 1)
                c2.text = element["inventory"]
                c2.text_frame.paragraphs[0].font.size = Pt(measures["cell_font"])
                table.rows[i].height = Cm(0.5)
                # Cell Color
                cell1 = table.cell(i, 0)
                cell2 = table.cell(i, 1)
                cell1.fill.solid()
                cell1.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell2.fill.solid()
                cell2.fill.fore_color.rgb = RGBColor(255, 255, 255)

            table.columns[0].width = Cm(3.8)
            table.columns[1].width = Cm(2.2)
        except Exception as e:
            raise SystemExit("Error: ", e)

    def create_pptx(self, ref_list: list[ProductData]):
        for idx, ref_data in enumerate(ref_list):
            self.create_title(ref_data["title"], idx, ref=ref_data["ref"])
            if "subtitle" in ref_data:
                self.create_subtitle(ref_data["subtitle"], idx)

            self.create_description(ref_data["description"], idx)
            self.create_img(ref_data["image"], idx)
            self.create_quantity_table(idx)
            self.create_inventory_table(ref_data["color_inventory"], idx)

    def save(self, path: str):
        self.prs.save(path)
